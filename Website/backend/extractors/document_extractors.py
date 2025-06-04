"""
Document text extraction utilities for various file formats.
Provides extraction functions for PDF, Word, and PowerPoint files.
"""

import os

# Functions for extracting text from different file types
def extract_text_from_pdf(file_path: str) -> str:
    """Extract text from PDF files."""
    try:
        # Try primary method with PyPDF
        try:
            from pypdf import PdfReader
            reader = PdfReader(file_path)
            text_parts = []
            for page in reader.pages:
                page_text = page.extract_text()
                if page_text and page_text.strip():
                    text_parts.append(page_text.strip())
            
            if text_parts:
                return "\n\n".join(text_parts)
            else:
                raise Exception("No text extracted using PyPDF")
        
        except Exception as pypdf_error:
            print(f"PyPDF error: {str(pypdf_error)}")
            
            # Fallback: Try pdfplumber if available
            try:
                import subprocess
                subprocess.check_call(["pip", "install", "pdfplumber"], stdout=subprocess.DEVNULL)
                
                import pdfplumber
                text_parts = []
                
                with pdfplumber.open(file_path) as pdf:
                    for page in pdf.pages:
                        page_text = page.extract_text() or ""
                        if page_text.strip():
                            text_parts.append(page_text.strip())
                
                if text_parts:
                    return "\n\n".join(text_parts)
                else:
                    raise Exception("No text extracted using pdfplumber")
                    
            except Exception as plumber_error:
                print(f"pdfplumber error: {str(plumber_error)}")
                
                # Last resort: try pdf2text if available
                try:
                    import subprocess
                    subprocess.check_call(["pip", "install", "pdftotext"], stdout=subprocess.DEVNULL)
                    
                    import pdftotext
                    with open(file_path, "rb") as f:
                        pdf = pdftotext.PDF(f)
                    
                    text_parts = [page.strip() for page in pdf if page.strip()]
                    if text_parts:
                        return "\n\n".join(text_parts)
                    else:
                        raise Exception("No text extracted using pdftotext")
                        
                except Exception as pdf2text_error:
                    # Final attempt with a basic approach
                    try:
                        import re
                        with open(file_path, 'rb') as f:
                            content = f.read()
                            # Extract any readable text
                            text = ''.join(chr(b) if 32 <= b < 127 else ' ' for b in content)
                            # Clean up the text
                            text = re.sub(r'\s+', ' ', text).strip()
                            # Filter meaningful parts
                            parts = [p for p in re.split(r'[.!?;]\s+', text) if len(p.strip()) > 10]
                            if parts:
                                return "\n\n".join(parts)
                            else:
                                return "Could not extract meaningful text from this PDF. It may be scanned or protected."
                    except Exception:
                        return "Failed to extract text from this PDF. It may be encrypted, scanned, or damaged."
    except Exception as e:
        return f"Error extracting text from PDF: {str(e)}"


def extract_text_from_word(file_path: str) -> str:
    """Extract text from Word documents."""
    try:
        # Try python-docx first (free and no watermark)
        try:
            import subprocess
            subprocess.check_call(["pip", "install", "python-docx"], stdout=subprocess.DEVNULL)
            import docx
            
            doc = docx.Document(file_path)
            paragraphs = []
            
            # Extract text from paragraphs without numbering them
            for para in doc.paragraphs:
                text = para.text.strip()
                if text:
                    paragraphs.append(text)
            
            # Also get text from tables
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        text = cell.text.strip()
                        if text:
                            paragraphs.append(text)
            
            if paragraphs:
                # Join paragraphs with double newlines without adding paragraph numbers
                return "\n\n".join(paragraphs)
            else:
                raise Exception("No text extracted using python-docx")
                
        except Exception as docx_error:
            print(f"python-docx error: {str(docx_error)}")
            
            # Fallback to docx2txt (another free library)
            try:
                import subprocess
                subprocess.check_call(["pip", "install", "docx2txt"], stdout=subprocess.DEVNULL)
                import docx2txt
                
                text = docx2txt.process(file_path)
                if text.strip():
                    # Process text to make it cleaner
                    lines = [line.strip() for line in text.split('\n') if line.strip()]
                    return "\n\n".join(lines)
                else:
                    raise Exception("No text extracted using docx2txt")
            
            except Exception as docx2txt_error:
                print(f"docx2txt error: {str(docx2txt_error)}")
                
                # Last resort: basic text extraction 
                try:
                    with open(file_path, 'rb') as f:
                        content = f.read()
                        # Extract any readable text (this is crude)
                        text = ''.join(chr(b) if 32 <= b < 127 else ' ' for b in content)
                        # Clean up the text by removing consecutive spaces
                        import re
                        text = re.sub(r'\s+', ' ', text).strip()
                        # Split by periods, commas, or other likely sentence boundaries
                        sentences = re.split(r'[.!?;]\s+', text)
                        # Filter out very short or likely meaningless sentences
                        sentences = [s.strip() for s in sentences if len(s.strip()) > 10]
                        return "\n\n".join(sentences)
                except Exception as basic_error:
                    return f"Failed to extract text from Word document using multiple methods. Try converting to text or PDF format."
    
    except Exception as e:
        return f"Error extracting text from Word document: {str(e)}"


def extract_text_from_ppt(file_path: str) -> str:
    """Extract text from PowerPoint presentations."""
    try:
        # Try python-pptx first (free library)
        try:
            import subprocess
            subprocess.check_call(["pip", "install", "python-pptx"], stdout=subprocess.DEVNULL)
            from pptx import Presentation
            
            prs = Presentation(file_path)
            slide_texts = []
            
            for slide in prs.slides:
                texts = []
                # Get text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        texts.append(shape.text.strip())
                        
                # Join all text from this slide
                if texts:
                    slide_texts.append("\n".join(texts))
            
            if slide_texts:
                return "\n\n".join(slide_texts)
            else:
                raise Exception("No text extracted using python-pptx")
                
        except Exception as pptx_error:
            print(f"python-pptx error: {str(pptx_error)}")
            
            # Basic text extraction as last resort
            try:
                with open(file_path, 'rb') as f:
                    content = f.read()
                    # Extract any readable text
                    import re
                    text = ''.join(chr(b) if 32 <= b < 127 else ' ' for b in content)
                    # Clean up the text
                    text = re.sub(r'\s+', ' ', text).strip()
                    # Filter meaningful parts
                    parts = [p for p in re.split(r'[.!?;]\s+', text) if len(p.strip()) > 10]
                    if parts:
                        return "\n\n".join(parts)
                    else:
                        return "Could not extract meaningful text from this presentation."
            except Exception:
                return "Failed to extract text from this PowerPoint file."
    except Exception as e:
        return f"Error extracting text from PowerPoint: {str(e)}" 